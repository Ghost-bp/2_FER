"""
生成 Phase3-组会汇报.pptx (基于旧 PPTX 模板改造)
策略: 复制旧 PPTX, 逐 slide 替换内容
"""
import copy, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

SRC = r"d:\MM_Emotion\2_FER\experiments\results\阶段二\Phase2-组会汇报.pptx"
DST = r"d:\MM_Emotion\2_FER\experiments\results\阶段三\Phase3-组会汇报.pptx"
CHART_DIR = r"d:\MM_Emotion\2_FER\experiments\results\阶段三\charts"

# ── 复制 PPTX ─────────────────────────────────────
prs = Presentation(SRC)

# ── 辅助函数 ───────────────────────────────────────
def clear_slide(slide, keep_indices=None):
    """删除 slide 中除 keep_indices 外的所有 shape"""
    keep = set(keep_indices or [])
    shapes_to_remove = []
    for i, shape in enumerate(slide.shapes):
        if i not in keep:
            shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

def set_text(shape, text, font_size=None, bold=None, color=None, alignment=None):
    """设置 shape 的第一段文字"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if font_size:
        p.font.size = Pt(font_size)
    if bold is not None:
        p.font.bold = bold
    if color:
        p.font.color.rgb = RGBColor(*color)
    if alignment:
        p.alignment = alignment

def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = RGBColor(*color)
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=11,
                   title_line=None, title_size=14, color=None):
    """添加多行文本框"""
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    if title_line:
        p = tf.paragraphs[0]
        p.text = title_line
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*(color or (0x2c, 0x7b, 0xb6)))
        p.space_after = Pt(8)

    for i, line in enumerate(lines):
        if title_line:
            p = tf.add_paragraph()
        elif i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(*(color or (0x33, 0x33, 0x33)))
        p.space_after = Pt(4)
    return txBox

def insert_chart(slide, left, top, width, height, chart_name):
    """插入 PNG 图表"""
    path = os.path.join(CHART_DIR, chart_name)
    if os.path.exists(path):
        return slide.shapes.add_picture(path, Emu(left), Emu(top), Emu(width), Emu(height))
    print(f"  WARN: chart not found: {path}")
    return None

def add_slide_title(slide, title, subtitle=None):
    """添加标准标题栏"""
    # 顶部分隔线
    line = slide.shapes.add_shape(
        1, Emu(0), Emu(0), Emu(12191695), Emu(73152))  # MSO_SHAPE.RECTANGLE
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x2c, 0x7b, 0xb6)
    line.line.fill.background()

    add_text_box(slide, 731520, 182880, 10698480, 548640, title,
                 font_size=28, bold=True, color=(0x2c, 0x7b, 0xb6))
    if subtitle:
        add_text_box(slide, 731520, 685800, 10698480, 365760, subtitle,
                     font_size=14, color=(0x66, 0x66, 0x66))

# ── Slide 0 (封面) ─────────────────────────────────
slide = prs.slides[0]
# Keep rectangles (0, 4), modify text boxes (1, 2, 3)
set_text(slide.shapes[1], "Phase 3: CAFE 全参数 + 结构优化", font_size=36, bold=True, color=(0xFF, 0xFF, 0xFF))
set_text(slide.shapes[2], "KMU-FED 6类表情识别 · KFold-10 受试者独立 · 21 实验", font_size=16, color=(0xCC, 0xDD, 0xFF))
set_text(slide.shapes[3], "2026-07-21", font_size=14, color=(0xAA, 0xBB, 0xDD))

# 背景色
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x1a, 0x2a, 0x3a)
# 标题文字颜色 (白色)
for s in [slide.shapes[1], slide.shapes[2], slide.shapes[3]]:
    for p in s.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide 1 (实验框架 L1-L7) ───────────────────────
slide = prs.slides[1]
clear_slide(slide, keep_indices=[])  # 全部重建
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "实验框架: 7 层分层优化",
                "控制变量 · 层级递进 · 早停自动终止 (patience=10)")

layers = [
    ("L1", "Backbone\n+ 预训练", "RN18/34/50\nMSCeleb/ImageNet", "5 实验"),
    ("L2", "优化器\n+ 调度器", "Adam/AdamW\nExp/Cos/Plateau", "4 实验"),
    ("L3", "学习率\n+ 权重衰减", "lr 5e-4~1e-4\nwd 1e-4~0", "5 实验"),
    ("L4", "数据增强", "Geometric\nColorJitter", "3 实验"),
    ("L5", "损失权重\n+ 正则化", "div=3/5/7\nLS/Dropout", "5 实验"),
    ("L6", "训练配置", "bs=32/64\nep=60/80", "2 实验"),
    ("L7", "最终验证", "3 种子统计\n×10Fold", "3 实验"),
]

COLORS_L = [(0x2c,0x7b,0xb6), (0xfd,0xae,0x61), (0xd7,0x19,0x1c), (0x1b,0x9e,0x77),
            (0x75,0x70,0xb3), (0xe7,0x29,0x8a), (0x66,0xa6,0x1e)]

card_w, card_h = 1500000, 2300000
start_x = 500000
gap = 100000
y = 1200000

for i, (label, title, detail, count) in enumerate(layers):
    x = start_x + i * (card_w + gap)
    # 卡片背景
    card = slide.shapes.add_shape(
        5, Emu(x), Emu(y), Emu(card_w), Emu(card_h))  # Rounded rectangle
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)
    card.line.color.rgb = RGBColor(*COLORS_L[i])
    card.line.width = Pt(1.5)

    # 标签圆圈
    circle = slide.shapes.add_shape(
        9, Emu(x + 50000), Emu(y + 80000), Emu(250000), Emu(250000))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(*COLORS_L[i])
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = label
    ctf.paragraphs[0].font.size = Pt(16)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + 350000, y + 60000, card_w - 400000, 350000,
                 title, font_size=13, bold=True, color=(0x33,0x33,0x33))
    add_text_box(slide, x + 100000, y + 450000, card_w - 200000, 1400000,
                 detail, font_size=10, color=(0x66,0x66,0x66))

    # 实验数标签
    cnt_badge = slide.shapes.add_shape(
        5, Emu(x + card_w - 750000), Emu(y + card_h - 380000), Emu(650000), Emu(280000))
    cnt_badge.fill.solid()
    cnt_badge.fill.fore_color.rgb = RGBColor(*COLORS_L[i])
    cnt_badge.line.fill.background()
    ctf2 = cnt_badge.text_frame
    ctf2.paragraphs[0].text = count
    ctf2.paragraphs[0].font.size = Pt(11)
    ctf2.paragraphs[0].font.bold = True
    ctf2.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    ctf2.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(slide, 500000, 3700000, 11000000, 400000,
             "策略: 参数间有依赖 → 串行, 无依赖 → 并行 | 每层锁定最优 → 进入下一层",
             font_size=11, color=(0x99,0x99,0x99))

print("  Slide 2 done")

# ── Slide 2 (优化路径总表) ─────────────────────────
slide = prs.slides[2]
clear_slide(slide)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "优化路径总表: 68.80% → 92.69%",
                "KFold-10 · 受试者独立 · 每步 +1.5~11pp")

# 表格
rows_data = [
    ["步骤", "Backbone", "预训练", "调度器", "增强", "损失权重", "Acc", "Std", "提升"],
    ["baseline", "RN18", "MSCeleb", "Exp", "HFlip+Erase", "div=5", "68.80%", "18.23%", "--"],
    ["+ImageNet", "RN18", "ImageNet", "Exp", "HFlip+Erase", "div=5", "79.84%", "15.63%", "+11.0"],
    ["+ResNet-34", "RN34", "ImageNet", "Exp", "HFlip+Erase", "div=5", "86.64%", "11.92%", "+6.8"],
    ["+ResNet-50", "RN50", "ImageNet", "Exp", "HFlip+Erase", "div=5", "88.42%", "11.12%", "+1.8"],
    ["+Cosine", "RN50", "ImageNet", "Cosine", "HFlip+Erase", "div=5", "89.24%", "8.97%", "+0.8"],
    ["+Geo增强", "RN50", "ImageNet", "Cosine", "+Geometric", "div=5", "90.74%", "10.75%", "+1.5"],
    ["+div=3", "RN50", "ImageNet", "Cosine", "+Geometric", "div=3", "92.22%", "8.90%", "+1.5"],
    ["+ep=80", "RN50", "ImageNet", "Cosine", "+Geometric", "div=3", "92.69%", "7.39%", "+0.5"],
]

nrows, ncols = len(rows_data), len(rows_data[0])
tbl_left = 400000
tbl_top = 1100000
tbl_w = 11400000
tbl_h = 5400000

table_shape = slide.shapes.add_table(nrows, ncols, Emu(tbl_left), Emu(tbl_top), Emu(tbl_w), Emu(tbl_h))
table = table_shape.table

# 列宽
col_widths = [1100000, 900000, 900000, 900000, 1200000, 1000000, 950000, 950000, 800000]
for i, w in enumerate(col_widths):
    table.columns[i].width = Emu(w)

for r in range(nrows):
    for c in range(ncols):
        cell = table.cell(r, c)
        cell.text = rows_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif r == nrows - 1:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xd7, 0x19, 0x1c)

        # 表头颜色
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2c, 0x7b, 0xb6)
        elif r == nrows - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)

add_text_box(slide, 400000, 6600000, 11000000, 250000,
             "总计: +23.89pp (vs 修复后基线) | 最终 3 种子均值: 90.63%",
             font_size=10, color=(0x99,0x99,0x99))

print("  Slide 3 done")

# ── Slide 3 (优化路径图 + Backbone分析) ────────────
slide = prs.slides[3]
clear_slide(slide)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "L1: Backbone + 预训练 (5 实验)",
                "RN18/34/50 × MSCeleb/ImageNet · Adam+ExpLR · HFlip+Erasing")

insert_chart(slide, 400000, 1100000, 7200000, 4500000, "optimization-path.png")

add_multi_text(slide, 7900000, 1100000, 3900000, 4500000, [
    "ImageNet 预训练优于 MSCeleb +11.0pp",
    "→ 通用视觉特征比人脸专有特征",
    "   更适合表情识别",
    "",
    "ResNet 深度越大越好",
    "→ RN34 > RN18 (+6.8pp)",
    "→ RN50 > RN34 (+1.8pp)",
    "→ 边际收益递减",
    "",
    "RN50 + ImageNet = 88.42%",
    "vs baseline 68.80% (+19.6pp)",
    "锁定进入 L2",
], font_size=11, title_line="关键发现", title_size=16, color=(0x2c,0x7b,0xb6))

print("  Slide 4 done")

# ── Slide 4 (L2 优化器+调度器) ─────────────────────
slide = prs.slides[4]
clear_slide(slide)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "L2: 优化器 + 调度器 (4 实验)",
                "RN50+ImageNet · lr=2e-4 · wd=1e-4")

insert_chart(slide, 400000, 1100000, 7200000, 4500000, "layer-comparison.png")

add_multi_text(slide, 7900000, 1100000, 3900000, 4500000, [
    "Adam > AdamW (+3.6pp)",
    "→ 耦合 Weight Decay 可能",
    "   更适合小数据集",
    "",
    "Cosine 调度器最优",
    "→ 89.24% vs Exp 88.42%",
    "→ 标准差从 11.12→8.97",
    "→ 平滑衰减均衡各 fold",
    "",
    "Plateau 表现最差 (85.69%)",
    "→ reduce-on-plateau 过于保守",
    "锁定: Adam + Cosine",
], font_size=11, title_line="关键发现", title_size=16, color=(0x2c,0x7b,0xb6))

print("  Slide 5 done")

# ── Slide 5 (L4+L5 增强+损失) ──────────────────────
slide = prs.slides[5]
clear_slide(slide)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "L4-L5: 数据增强 + 损失权重 + 正则化",
                "RN50+Adam+Cosine · 8 实验")

# 左: acc-loss-summary
insert_chart(slide, 400000, 1100000, 5500000, 3200000, "acc-loss-summary.png")
# 右: best-loss-curves
insert_chart(slide, 6200000, 1100000, 5500000, 3200000, "best-loss-curves.png")

add_multi_text(slide, 400000, 4500000, 11000000, 2000000, [
    "L4 数据增强: Geometric +1.5pp (Rotation 15deg + Affine 0.1)",
    "            ColorJitter + Normalize 致命 (23.52%), 颜色增强在归一化后失效",
    "",
    "L5 损失权重: div=3 > div=7 > div=5, 降低多样性约束利于主任务 (+1.5pp)",
    "             LabelSmoothing 有害 (-1.8pp), Dropout=0.3 中性 (-0.05pp)",
    "             小数据集 FER 场景, 正则化整体不推荐",
], font_size=11, title_line="关键发现", title_size=14, color=(0x2c,0x7b,0xb6))

print("  Slide 6 done")

# ── Slide 6 (L7 最终结果) ──────────────────────────
slide = prs.slides[6]
clear_slide(slide)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "最终验证: 3 种子 × 10 Fold",
                "最优配置: RN50+ImageNet+Adam+Cosine+lr2e-4+wd1e-4+Geo+div3+ep80")

insert_chart(slide, 400000, 1100000, 5800000, 4200000, "final-3seeds.png")
insert_chart(slide, 6400000, 1100000, 5500000, 4200000, "baseline-vs-best-folds.png")

add_multi_text(slide, 400000, 5500000, 11000000, 1200000, [
    "种子 42: 92.69% ± 7.39%  |  种子 123: 89.74% ± 11.92%  |  种子 456: 89.46% ± 12.25%",
    "30-Fold 均值: 90.63%  |  vs 修复基线 68.80%: +21.83pp  |  vs 师兄原版 74.38%: +16.25pp",
], font_size=12, title_line="最终结果 (3 种子均值 = 90.63%)", title_size=16, color=(0xd7,0x19,0x1c))

print("  Slide 7 done")

# ── Slide 7 (综合分析) ─────────────────────────────
slide = prs.slides[7]
clear_slide(slide)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "Acc + Loss 综合分析",
                "训练 vs 验证: 过拟合程度 · 收敛速度 · Loss 变化趋势")

insert_chart(slide, 400000, 1100000, 11400000, 5300000, "best-loss-curves.png")

add_text_box(slide, 400000, 6550000, 11000000, 250000,
             "Train Loss: CE + 3*Div + 1.5*Mask 复合损失  |  Val Loss ≈ 1.3~1.5  |  过拟合可控 (Val Acc 持续提升)",
             font_size=10, color=(0x99,0x99,0x99))

print("  Slide 8 done")

# ── Slide 8 (各维度总结表) ─────────────────────────
slide = prs.slides[8]
clear_slide(slide)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "各维度最优选择总结",
                "基于 21 实验 · 控制变量逐层验证")

summary_data = [
    ["维度", "最优选择", "次优选择", "关键说明"],
    ["Backbone", "ResNet-50", "ResNet-34", "+1.8pp vs RN34, 最多3个fold达100%"],
    ["预训练", "ImageNet", "MSCeleb (仅RN18)", "+11.0pp, 通用特征>人脸专有特征"],
    ["优化器", "Adam", "AdamW", "+3.6pp, 耦合WD更适合小数据集"],
    ["调度器", "CosineAnnealing", "Exponential", "+0.8pp, std减半(8.97%)"],
    ["学习率", "2e-4", "1e-4/5e-4", "默认值恰是最优, 过大/过小均降~6pp"],
    ["权重衰减", "1e-4", "1e-5", "默认值最优, wd=0 降3pp"],
    ["数据增强", "Geometric", "HFlip+Erase(基础)", "+1.5pp, ColorJitter禁用(23.5%)"],
    ["损失权重(div)", "3.0", "7.0", "+1.5pp, 弱多样性约束利于主任务"],
    ["训练轮数", "80 (早停)", "60", "+0.5pp, std降至7.39%"],
    ["正则化", "不使用", "Dropout=0.3", "小数据集FER场景, 全部有害/中性"],
    ["最终配置: 92.69%", "", "", "RN50+ImageNet+Adam+Cosine+Geo+div3+ep80"],
]

nrows2 = len(summary_data)
ncols2 = len(summary_data[0])
table_shape2 = slide.shapes.add_table(nrows2, ncols2, Emu(400000), Emu(1100000), Emu(11400000), Emu(5400000))
table2 = table_shape2.table

col_widths2 = [1500000, 2300000, 2300000, 5300000]
for i, w in enumerate(col_widths2):
    table2.columns[i].width = Emu(w)

for r in range(nrows2):
    for c in range(ncols2):
        cell = table2.cell(r, c)
        cell.text = summary_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER if c < 3 else PP_ALIGN.LEFT
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            elif r == nrows2 - 1:
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xd7, 0x19, 0x1c)
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2c, 0x7b, 0xb6)
        elif r == nrows2 - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF8, 0xFC)

add_text_box(slide, 400000, 6580000, 11000000, 250000,
             "vs 修复基线(68.80%): +23.89pp  |  vs 师兄原版(74.38%): +18.31pp  |  3种子均值: 90.63%",
             font_size=10, color=(0x99,0x99,0x99))

print("  Slide 9 done")

# ── Slide 9 (后续工作) ─────────────────────────────
slide = prs.slides[9]
clear_slide(slide)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
add_slide_title(slide, "后续工作与展望")

next_steps = [
    ("1", "论文撰写", "基于最优配置(92.69%)撰写实验部分, 对比 SOTA 方法"),
    ("2", "跨数据集验证", "在 RAF-DB, AffectNet, FER2013 上验证泛化能力"),
    ("3", "更大 Backbone", "ResNet-101, ConvNeXt, ViT 等更强 backbone"),
    ("4", "CLIP 变体", "CLIP ViT-L/14, SigLIP 等更大的视觉编码器"),
    ("5", "多模态融合", "结合 CLIP 文本分支做 zero-shot / few-shot FER"),
    ("6", "消融实验补全", "Geo 参数调优, div 细化搜索, MixUp 重测"),
]

for i, (num, title, desc) in enumerate(next_steps):
    y = 1200000 + i * 800000
    circle = slide.shapes.add_shape(9, Emu(600000), Emu(y), Emu(400000), Emu(400000))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x2c, 0x7b, 0xb6)
    circle.line.fill.background()
    ctf = circle.text_frame
    ctf.paragraphs[0].text = num
    ctf.paragraphs[0].font.size = Pt(18)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, 1200000, y + 20000, 3500000, 350000, title,
                 font_size=16, bold=True, color=(0x33,0x33,0x33))
    add_text_box(slide, 1200000, y + 370000, 9000000, 350000, desc,
                 font_size=12, color=(0x66,0x66,0x66))

print("  Slide 10 done")

# ── 保存 ──────────────────────────────────────────
os.makedirs(os.path.dirname(DST), exist_ok=True)
prs.save(DST)
print(f"\nDONE: {DST}")
