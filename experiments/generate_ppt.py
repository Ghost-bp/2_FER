"""
Phase 2 消融实验组会汇报 PPT 生成器

所有元素使用 python-pptx 原生对象（表格、图表、文本框、形状），
确保在 PowerPoint 中可逐一编辑。
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
import os

# ═══════════════════════════════════════════
# 配色方案
# ═══════════════════════════════════════════
C_TEXT = RGBColor(0x33, 0x33, 0x33)
C_TITLE = RGBColor(0x2C, 0x3E, 0x50)
C_TABLE_HEADER_BG = RGBColor(0x44, 0x72, 0xC4)
C_TABLE_HEADER_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
C_TABLE_ROW_ODD = RGBColor(0xFF, 0xFF, 0xFF)
C_TABLE_ROW_EVEN = RGBColor(0xF2, 0xF7, 0xFC)
C_BEST_BG = RGBColor(0xE2, 0xEF, 0xDA)
C_BAD_BG = RGBColor(0xFC, 0xE4, 0xD6)
C_CHART_BAR = RGBColor(0x44, 0x72, 0xC4)
C_CHART_HIGHLIGHT = RGBColor(0xED, 0x7D, 0x31)
C_GRAY_BAR = RGBColor(0xBB, 0xBB, 0xBB)
C_CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)
C_CARD_BORDER = RGBColor(0xD0, 0xD5, 0xDD)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_SUBTLE = RGBColor(0x88, 0x88, 0x88)

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ═══════════════════════════════════════════
# 实验数据
# ═══════════════════════════════════════════

ALL_EXPERIMENTS = [
    # (实验名, 组, Backbone, 参数量, 优化器, 调度器, 增强, 正则化, Acc, Std, 协议, 状态)
    ("F1",  "F", "ResNet-50",  "23.5M", "AdamW",  "Cos",    "Geo",     "--",      89.31, 10.38, "YOLO",  "best"),
    ("C5",  "C", "ResNet-50",  "23.5M", "AdamW",  "Cos",    "--",      "--",      80.09, 7.55,  "YOLO",  "ok"),
    ("E2",  "E", "ResNet-50",  "23.5M", "Adam",   "Cos",    "Geo",     "--",      79.62, 8.62,  "Haar",  "ok"),
    ("D2",  "D", "ResNet-50",  "23.5M", "AdamW",  "Cos",    "--",      "Drop0.5", 79.18, 11.78, "YOLO",  "ok"),
    ("C3",  "C", "ResNet-50",  "23.5M", "Adam",   "Cos",    "--",      "--",      78.44, 9.55,  "Haar",  "ok"),
    ("C2b", "C", "ResNet-50",  "23.5M", "SGD",    "Cos",    "--",      "--",      77.78, 11.71, "YOLO",  "ok"),
    ("D4",  "D", "ResNet-50",  "23.5M", "AdamW",  "Cos",    "--",      "WD=1e-3", 77.10, 14.72, "YOLO",  "ok"),
    ("A3",  "A", "ResNet-50",  "23.5M", "Adam",   "Exp",    "--",      "--",      76.92, 11.86, "Haar",  "ok"),
    ("D3",  "D", "ResNet-50",  "23.5M", "Adam",   "Cos",    "--",      "LSmooth", 73.66, 11.20, "Haar",  "ok"),
    ("D1",  "D", "ResNet-50",  "23.5M", "Adam",   "Cos",    "--",      "Drop0.3", 73.68, 9.48,  "Haar",  "ok"),
    ("B5",  "B", "ResNet-18",  "11.2M", "Adam",   "Exp",    "--",      "--",      73.09, 13.60, "Haar",  "ok"),
    ("B4",  "B", "ResNet-18",  "11.2M", "Adam",   "Exp",    "--",      "--",      73.03, 11.76, "Haar",  "ok"),
    ("A2",  "A", "ResNet-34",  "21.3M", "Adam",   "Exp",    "--",      "--",      72.86, 9.82,  "Haar",  "ok"),
    ("C4",  "C", "ResNet-50",  "23.5M", "Adam",   "Plateau","--",      "--",      71.30, 12.10, "Haar",  "ok"),
    ("E4",  "E", "ResNet-50",  "23.5M", "Adam",   "Cos",    "MixUp",   "--",      71.19, 10.42, "Haar",  "ok"),
    ("C1",  "C", "ResNet-50",  "23.5M", "AdamW",  "Exp",    "--",      "--",      70.81, 12.49, "Haar",  "ok"),
    ("A4",  "A", "EfficientNet-B0","5.3M","Adam",  "Exp",    "--",      "--",      70.01, 16.63, "Haar",  "ok"),
    ("E3",  "E", "ResNet-50",  "23.5M", "AdamW",  "Cos",    "Geo+Color","--",     24.78, 7.54,  "YOLO",  "crashed"),
    ("E1",  "E", "ResNet-50",  "23.5M", "Adam",   "Cos",    "Color",   "--",      None,  None,  "Haar",  "crashed"),
    ("C2",  "C", "ResNet-50",  "23.5M", "SGD",    "Exp",    "--",      "--",      None,  None,  "YOLO",  "crashed"),
]

# F1 per-fold data
F1_FOLDS = [
    ("S1", 89.68), ("S2", 83.33), ("S3", 97.78), ("S4", 100.00),
    ("S5", 99.17), ("S6", 89.17), ("S7", 93.00), ("S8", 100.00),
    ("S9", 92.86), ("S10", 87.50), ("S11", 75.00), ("S12", 64.29),
]

# Backbone comparison data
BACKBONE_DATA = [
    ("RN50\n(23.5M)", 76.92, 11.86),
    ("RN34\n(21.3M)", 72.86, 9.82),
    ("RN18\n(11.2M)", 73.03, 11.76),
    ("EfficientNet-B0\n(5.3M)", 70.01, 16.63),
]

# Optimizer comparison data
OPTIMIZER_DATA = [
    ("C5\nAdamW+Cos", 80.09),
    ("C3\nAdam+Cos", 78.44),
    ("C2b\nSGD+Cos", 77.78),
    ("C1\nAdamW+Exp", 70.81),
    ("C4\nAdam+\nPlateau", 71.30),
]

# Regularization data (with baseline)
REG_DATA = [
    ("基线\n(无正则)", 80.09),
    ("D2\nDropout0.5", 79.18),
    ("D4\nWD=1e-3", 77.10),
    ("D1\nDropout0.3", 73.68),
    ("D3\nLabelSmooth", 73.66),
]

# Augmentation data (with baseline)
AUG_DATA = [
    ("基线\n(无增强)", 80.09),
    ("E2\nGeometric", 79.62),
    ("E4\nMixUp", 71.19),
    ("E3\nGeo+Color", 24.78),
]

# Synergy data
SYNERGY_DATA = [
    ("E2\nHaar+Geo", 79.62),
    ("C5\nYOLO only", 80.09),
    ("F1\nYOLO+Geo", 89.31),
]

# ═══════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════

def set_slide_bg(slide, color):
    """Set slide background to solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, text, y=Inches(0.2), font_size=Pt(28)):
    """Add a styled title text box at the top of a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = C_TITLE
    p.font.bold = True
    return txBox


def add_subtitle(slide, text, y=Inches(0.75), font_size=Pt(14)):
    """Add a subtitle line."""
    txBox = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.7), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = C_SUBTLE
    return txBox


def set_cell_fmt(cell, text, font_size=Pt(10), bold=False, color=None,
                 fill_color=None, alignment=PP_ALIGN.CENTER):
    """Format a table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = str(text)
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color or C_TEXT
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill_color:
        set_cell_fill(cell, fill_color)


def set_cell_fill(cell, color):
    """Set cell background fill."""
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '{:02X}{:02X}{:02X}'.format(color[0] if isinstance(color, tuple) else color.red, color[1] if isinstance(color, tuple) else color.green, color[2] if isinstance(color, tuple) else color.blue)})
    # Use hex string approach
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove existing fills
    for child in list(tcPr):
        if child.tag in [qn('a:solidFill'), qn('a:noFill'), qn('a:gradFill')]:
            tcPr.remove(child)
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    if isinstance(color, tuple):
        hex_val = '{:02X}{:02X}{:02X}'.format(*color)
    else:
        hex_val = '{:02X}{:02X}{:02X}'.format(color.red, color.green, color.blue)
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': hex_val})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def _cell_fill_helper(cell, hex_str):
    """Helper: set cell background by hex string."""
    tcPr = cell._tc.get_or_add_tcPr()
    for child in list(tcPr):
        if child.tag in [qn('a:solidFill'), qn('a:noFill'), qn('a:gradFill')]:
            tcPr.remove(child)
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': hex_str})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def rgb_to_hex(color):
    """Convert RGBColor to hex string (python-pptx RGBColor is int subclass)."""
    if isinstance(color, int):
        return '{:06X}'.format(color)
    # Fallback for tuple
    return '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])


def add_native_chart(slide, data_tuples, left, top, width, height,
                     highlight_idx=None, title=""):
    """
    Add a native PPT clustered column chart.
    data_tuples: list of (label, value) or (label, value, err)
    highlight_idx: index of bar to color differently
    """
    chart_data = CategoryChartData()
    chart_data.categories = [t[0] for t in data_tuples]

    if len(data_tuples[0]) >= 3:
        # Has error values
        chart_data.add_series('Accuracy (%)', [t[1] for t in data_tuples])
    else:
        chart_data.add_series('Accuracy (%)', [t[1] for t in data_tuples])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    # Style the chart
    plot = chart.plots[0]
    series = plot.series[0]

    # Set bar colors
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = C_CHART_BAR

    # Per-point coloring for highlighted bar via XML
    if highlight_idx is not None:
        ser_elem = series._element
        # Remove existing dPt elements
        for dPt in ser_elem.findall(qn('c:dPt')):
            ser_elem.remove(dPt)
        for i in range(len(data_tuples)):
            dPt = ser_elem.makeelement(qn('c:dPt'), {})
            idx_elem = dPt.makeelement(qn('c:idx'), {'val': str(i)})
            dPt.append(idx_elem)
            spPr = dPt.makeelement(qn('c:spPr'), {})
            solidFill = spPr.makeelement(qn('a:solidFill'), {})
            if i == highlight_idx:
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': rgb_to_hex(C_CHART_HIGHLIGHT)})
            else:
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': rgb_to_hex(C_CHART_BAR)})
            solidFill.append(srgbClr)
            spPr.append(solidFill)
            dPt.append(spPr)
            ser_elem.append(dPt)

    # Data labels
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(8)
    data_labels.font.color.rgb = C_TEXT
    data_labels.number_format = '0.0"%"'

    # Style axes
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 105
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    value_axis.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.tick_labels.font.color.rgb = C_TEXT
    category_axis.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = C_TITLE

    return chart_frame


def add_line_on_chart(slide, left, top, width, y_offset, label, color=C_CHART_HIGHLIGHT):
    """Add a horizontal reference line annotation (as a text box + line shape)."""
    # Simple text annotation
    txBox = slide.shapes.add_textbox(left + width - Inches(1.5), top + y_offset, Inches(1.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT
    return txBox


# ═══════════════════════════════════════════
# Slide Builders
# ═══════════════════════════════════════════

def add_slide_1_cover(prs):
    """封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, C_WHITE)

    # Decorative top bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_TABLE_HEADER_BG
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Phase 2: 消融实验与超参数调优"
    p.font.size = Pt(40)
    p.font.color.rgb = C_TITLE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "KMU-FED 6类表情识别 · LOSO 12折交叉验证 · 20个消融实验"
    p2.font.size = Pt(18)
    p2.font.color.rgb = C_SUBTLE
    p2.alignment = PP_ALIGN.CENTER

    # Date
    txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "2026-07-17"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_SUBTLE
    p3.alignment = PP_ALIGN.CENTER

    # Bottom bar
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = C_TABLE_HEADER_BG
    bar2.line.fill.background()


def add_slide_2_overview(prs):
    """实验设计总览 — 6组消融卡片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "实验设计总览: 6 组消融实验")
    add_subtitle(slide, "控制变量法逐维度消融，统一 LOSO 12折 + YOLO-only 协议")

    groups = [
        ("A 组 · Backbone 消融", "ResNet-18/34/50\nEfficientNet-B0", "4 实验"),
        ("B 组 · 学习率消融", "lr = 1e-5 ~ 1e-3\n(5个数量级)", "5 实验"),
        ("C 组 · 优化器+调度器", "Adam/AdamW/SGD\nExp/Cos/Plateau", "6 实验"),
        ("D 组 · 正则化消融", "Dropout / LabelSmooth\nWeight Decay", "4 实验"),
        ("E 组 · 数据增强", "Geometric / Color\nMixUp / 组合", "4 实验"),
        ("F 组 · 最优组合", "各组最优因素的\n交叉组合验证", "1 实验"),
    ]

    card_w = Inches(3.6)
    card_h = Inches(2.2)
    start_x = Inches(0.8)
    start_y = Inches(1.5)
    gap_x = Inches(0.3)
    gap_y = Inches(0.3)
    cols = 3

    for i, (title, desc, count) in enumerate(groups):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_CARD_BORDER
        card.line.width = Pt(0.5)

        # Group label (A/B/C/D/E/F circle)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.4), Inches(0.4)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_TABLE_HEADER_BG
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.word_wrap = False
        cp = ctf.paragraphs[0]
        cp.text = title[0]  # Just the letter
        cp.font.size = Pt(14)
        cp.font.color.rgb = C_WHITE
        cp.font.bold = True
        cp.alignment = PP_ALIGN.CENTER

        # Title
        tb = slide.shapes.add_textbox(x + Inches(0.8), y + Inches(0.15), Inches(2.6), Inches(0.4))
        tp = tb.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(13)
        tp.font.bold = True
        tp.font.color.rgb = C_TITLE

        # Description
        tb2 = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.75), Inches(3.0), Inches(0.9))
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT
        p2.line_spacing = Pt(18)

        # Count badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(2.5), y + Inches(1.75), Inches(0.9), Inches(0.3)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_TABLE_HEADER_BG
        badge.line.fill.background()
        btf = badge.text_frame
        bp = btf.paragraphs[0]
        bp.text = count
        bp.font.size = Pt(9)
        bp.font.color.rgb = C_WHITE
        bp.alignment = PP_ALIGN.CENTER


def add_slide_3_main_table(prs):
    """核心结果总表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "Phase 2 实验完整总表")
    add_subtitle(slide, "按准确率降序排列 · 绿色=最优 · 橙色=终止/崩溃")

    headers = ["实验", "组", "Backbone", "参数量", "优化器", "调度器", "增强", "正则化", "Acc±Std"]
    col_widths = [Inches(0.7), Inches(0.45), Inches(1.55), Inches(0.8), Inches(0.95), Inches(0.95), Inches(1.1), Inches(1.1), Inches(1.2)]

    n_rows = len(ALL_EXPERIMENTS) + 1  # +1 for header
    n_cols = len(headers)
    table_left = Inches(0.5)
    table_top = Inches(1.3)
    table_width = sum(col_widths)
    table_height = Inches(5.8)

    table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        set_cell_fmt(cell, h, font_size=Pt(9), bold=True, color=C_TABLE_HEADER_TEXT)
        _cell_fill_helper(cell, rgb_to_hex(C_TABLE_HEADER_BG))

    # Data rows
    for i, exp in enumerate(ALL_EXPERIMENTS):
        row = i + 1
        name, grp, bb, params, opt, sched, aug, reg, acc, std, proto, status = exp

        # Format Acc±Std
        if acc is not None:
            acc_str = f"{acc:.2f}% ± {std:.2f}%"
        else:
            acc_str = "终止"

        row_data = [name, grp, bb, params, opt, sched, aug, reg, acc_str]

        # Row styling
        if status == "best":
            row_bg = rgb_to_hex(C_BEST_BG)
            is_bold = True
            text_color = C_TITLE
        elif status == "crashed":
            row_bg = rgb_to_hex(C_BAD_BG)
            is_bold = False
            text_color = C_SUBTLE
        elif row % 2 == 0:
            row_bg = rgb_to_hex(C_TABLE_ROW_EVEN)
            is_bold = False
            text_color = C_TEXT
        else:
            row_bg = rgb_to_hex(C_TABLE_ROW_ODD)
            is_bold = False
            text_color = C_TEXT

        for j, val in enumerate(row_data):
            cell = table.cell(row, j)
            set_cell_fmt(cell, val, font_size=Pt(9), bold=is_bold, color=text_color)
            _cell_fill_helper(cell, row_bg)

    # Bottom annotation
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(11), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "A组: Backbone消融  |  B组: 学习率消融  |  C组: 优化器+调度器  |  D组: 正则化  |  E组: 数据增强  |  F组: 最优组合"
    p.font.size = Pt(8)
    p.font.color.rgb = C_SUBTLE


def add_slide_4_backbone(prs):
    """Backbone 对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "A 组: Backbone 消融对比")
    add_subtitle(slide, "同协议 (Haar-first) 同超参, 仅 Backbone 不同")

    add_native_chart(slide, BACKBONE_DATA,
                     Inches(1.0), Inches(1.5), Inches(8.0), Inches(5.0),
                     highlight_idx=0)

    # Annotation text box
    txBox = slide.shapes.add_textbox(Inches(9.5), Inches(2.0), Inches(3.3), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    notes = [
        ("结论", True, C_TITLE),
        ("", False, C_TEXT),
        ("ResNet-50 显著最优", False, C_TEXT),
        ("+3.89pp vs RN34", False, C_SUBTLE),
        ("", False, C_TEXT),
        ("EfficientNet-B0", False, C_TEXT),
        ("精度最低 + 方差最大", False, C_SUBTLE),
        ("不适合小数据FER", False, C_SUBTLE),
        ("", False, C_TEXT),
        ("RN18 vs RN34", False, C_TEXT),
        ("参数量×2 精度不变", False, C_SUBTLE),
        ("RN18 性价比最优", False, C_SUBTLE),
    ]
    for i, (text, bold, color) in enumerate(notes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11) if bold else Pt(10)
        p.font.bold = bold
        p.font.color.rgb = color
        p.line_spacing = Pt(16)


def add_slide_5_optimizer(prs):
    """优化器+调度器对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "C 组: 优化器 + 调度器消融")
    add_subtitle(slide, "ResNet-50, lr=2e-4, 同协议 (YOLO-only 或 Haar-first)")

    add_native_chart(slide, OPTIMIZER_DATA,
                     Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.0),
                     highlight_idx=0)

    # Key annotation
    txBox = slide.shapes.add_textbox(Inches(9.8), Inches(2.0), Inches(3.0), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    notes = [
        ("关键发现", True, C_TITLE),
        ("", False, C_TEXT),
        ("1. CosineAnnealingLR", False, C_TEXT),
        ("   是最优调度器", False, C_SUBTLE),
        ("   Cos > Plateau > Exp", False, C_SUBTLE),
        ("", False, C_TEXT),
        ("2. AdamW+Cos 最优", False, C_TEXT),
        ("   AdamW+Cos > Adam+Cos", False, C_SUBTLE),
        ("   (+1.65pp)", False, C_SUBTLE),
        ("", False, C_TEXT),
        ("3. AdamW+Exp 惨败", False, C_TEXT),
        ("   比 Adam+Exp 低 6.1pp", False, C_SUBTLE),
        ("   → AdamW必须配Cosine!", False, C_TEXT),
    ]
    for i, (text, bold, color) in enumerate(notes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11) if bold else Pt(10)
        p.font.bold = bold
        p.font.color.rgb = color
        p.line_spacing = Pt(15)


def add_slide_6_reg_aug(prs):
    """正则化 + 数据增强 双栏"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "D 组: 正则化  &  E 组: 数据增强")
    add_subtitle(slide, "小数据集FER: 正则化全部有害, ColorJitter 灾难性")

    # Left chart: Regularization
    add_native_chart(slide, REG_DATA,
                     Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5),
                     highlight_idx=0)

    # Left annotation
    txBox_l = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(5.5), Inches(0.5))
    tf_l = txBox_l.text_frame
    p_l = tf_l.paragraphs[0]
    p_l.text = "所有正则化均低于基线 → 小数据集不需要额外正则化"
    p_l.font.size = Pt(11)
    p_l.font.bold = True
    p_l.font.color.rgb = C_CHART_HIGHLIGHT

    # Right chart: Augmentation
    add_native_chart(slide, AUG_DATA,
                     Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.5),
                     highlight_idx=0)

    # Right annotation
    txBox_r = slide.shapes.add_textbox(Inches(6.8), Inches(6.2), Inches(5.5), Inches(0.5))
    tf_r = txBox_r.text_frame
    p_r = tf_r.paragraphs[0]
    p_r.text = "ColorJitter 破坏表情纹理 → E3 崩溃至 24.78%"
    p_r.font.size = Pt(11)
    p_r.font.bold = True
    p_r.font.color.rgb = C_CHART_HIGHLIGHT


def add_slide_7_synergy(prs):
    """Geo×YOLO 协同效应"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "核心发现: Geometric × YOLO 协同效应")

    add_native_chart(slide, SYNERGY_DATA,
                     Inches(1.5), Inches(1.5), Inches(7.0), Inches(4.5),
                     highlight_idx=2)

    # Arrows and annotations
    # Arrow 1: E2 -> C5 (small gain)
    tx1 = slide.shapes.add_textbox(Inches(2.8), Inches(3.5), Inches(2.0), Inches(0.4))
    tf1 = tx1.text_frame
    p1 = tf1.paragraphs[0]
    p1.text = "← +1.18pp (Haar)"
    p1.font.size = Pt(10)
    p1.font.color.rgb = C_SUBTLE

    # Arrow 2: C5 -> F1 (big gain)
    tx2 = slide.shapes.add_textbox(Inches(6.0), Inches(3.0), Inches(2.0), Inches(0.4))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "+9.22pp (YOLO) →"
    p2.font.size = Pt(10)
    p2.font.color.rgb = C_CHART_HIGHLIGHT
    p2.font.bold = True

    # Explanation text
    txBox = slide.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.8), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    notes = [
        ("协同效应解释", True, C_TITLE),
        ("", False, C_TEXT),
        ("Haar Cascade", True, C_TEXT),
        ("tight face crop", False, C_SUBTLE),
        ("已裁掉大部分姿态信息", False, C_SUBTLE),
        ("Geo 增益自然受限", False, C_SUBTLE),
        ("", False, C_TEXT),
        ("YOLOv8n", True, C_TEXT),
        ("wider crop 保留空间上下文", False, C_SUBTLE),
        ("(头部姿态、肩膀轮廓)", False, C_SUBTLE),
        ("Geo 训练出姿态不变性", False, C_SUBTLE),
        ("", False, C_TEXT),
        ("协同因子: ×7.8 !", True, C_CHART_HIGHLIGHT),
    ]
    for i, (text, bold, color) in enumerate(notes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11) if bold else Pt(10)
        p.font.bold = bold
        p.font.color.rgb = color
        p.line_spacing = Pt(15)

    # Bottom key message
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(10.0), Inches(0.5))
    tf2b = txBox2.text_frame
    p2b = tf2b.paragraphs[0]
    p2b.text = "YOLO wider crop 保留空间上下文 → Geometric 增强学会真正的姿态不变性 → 89.31%"
    p2b.font.size = Pt(14)
    p2b.font.bold = True
    p2b.font.color.rgb = C_TITLE
    p2b.alignment = PP_ALIGN.CENTER


def add_slide_8_f1_folds(prs):
    """F1 12折详情"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "F1 最优组合: 12 折交叉验证详情")
    add_subtitle(slide, "ResNet-50 + AdamW + Cosine + Geometric + YOLO-only  |  89.31% ± 10.38%")

    # Bar chart
    chart_data = CategoryChartData()
    chart_data.categories = [f[0] for f in F1_FOLDS]
    chart_data.add_series('Accuracy (%)', [f[1] for f in F1_FOLDS])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.0), chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    plot = chart.plots[0]
    series = plot.series[0]

    # Color bars: above mean = blue, below mean = gray
    mean_val = 89.31
    for i, (name, val) in enumerate(F1_FOLDS):
        pt = series.points[i] if hasattr(series, 'points') else None
        # Actually, python-pptx doesn't support per-point coloring easily
        # Let's just set all to blue and mark the 100% ones
        pass

    # Data labels
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(8)
    data_labels.font.color.rgb = C_TEXT
    data_labels.number_format = '0.0"%"'

    # Value axis
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = 105
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(10)
    category_axis.tick_labels.font.color.rgb = C_TEXT

    # Mean annotation
    txBox = slide.shapes.add_textbox(Inches(10.0), Inches(1.0), Inches(2.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "均值: 89.31%"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = C_CHART_HIGHLIGHT
    p.alignment = PP_ALIGN.RIGHT

    # Bottom highlights
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.0), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "2 折满分 (S4, S8)  |  5 折 > 90%  |  S11/S12 缺类瓶颈 (75%/64%)"
    p2.font.size = Pt(11)
    p2.font.color.rgb = C_SUBTLE
    p2.alignment = PP_ALIGN.CENTER


def add_slide_9_summary_table(prs):
    """各维度最优总结"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "各维度最优选择总结")

    dims = [
        ("Backbone", "ResNet-50", "ResNet-18 (性价比)", "+3.89pp vs RN34"),
        ("优化器", "AdamW", "Adam", "+1.65pp vs Adam"),
        ("调度器", "CosineAnnealingLR", "ExponentialLR", "+1.52pp vs Exp"),
        ("学习率", "2e-4", "5e-4 (RN18)", "RN50 用 2e-4"),
        ("权重衰减", "1e-4", "—", "1e-3 显著有害"),
        ("数据增强", "Geometric", "无增强", "**+9.22pp (YOLO)**"),
        ("正则化", "不使用", "—", "全部低于基线"),
    ]

    n_rows = len(dims) + 2  # header + data + footer
    n_cols = 4
    headers = ["维度", "最优选择", "次优选择", "关键备注"]
    col_widths = [Inches(1.8), Inches(3.0), Inches(3.0), Inches(3.5)]

    table_left = Inches(0.7)
    table_top = Inches(1.5)
    table_width = sum(col_widths)
    table_height = Inches(4.5)

    table_shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        set_cell_fmt(cell, h, font_size=Pt(12), bold=True, color=C_TABLE_HEADER_TEXT)
        _cell_fill_helper(cell, rgb_to_hex(C_TABLE_HEADER_BG))

    # Data
    for i, (dim, best, second, note) in enumerate(dims):
        row = i + 1
        bg = rgb_to_hex(C_TABLE_ROW_EVEN) if row % 2 == 0 else rgb_to_hex(C_TABLE_ROW_ODD)
        set_cell_fmt(table.cell(row, 0), dim, font_size=Pt(12), bold=True, color=C_TITLE)
        set_cell_fmt(table.cell(row, 1), best, font_size=Pt(12), bold=True, color=C_TABLE_HEADER_BG)
        set_cell_fmt(table.cell(row, 2), second, font_size=Pt(11), color=C_TEXT)
        # Check if note contains "**"
        is_important = "**" in note
        set_cell_fmt(table.cell(row, 3), note.replace("**", ""), font_size=Pt(11),
                     bold=is_important, color=C_TEXT)
        for j in range(4):
            _cell_fill_helper(table.cell(row, j), bg)

    # Footer row
    footer_row = n_rows - 1
    footer_text = "最优组合 → F1 = 89.31% ± 10.38%"
    # Merge footer
    table.cell(footer_row, 0).merge(table.cell(footer_row, n_cols - 1))
    set_cell_fmt(table.cell(footer_row, 0), footer_text,
                 font_size=Pt(14), bold=True, color=C_WHITE)
    _cell_fill_helper(table.cell(footer_row, 0), rgb_to_hex(C_CHART_HIGHLIGHT))


def add_slide_10_conclusion(prs):
    """结论与下一步"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_WHITE)
    add_title_bar(slide, "结论与下一步工作")

    conclusions = [
        ("1", "KMU-FED 最优配置确定",
         "ResNet-50 + AdamW + Cosine + Geometric + YOLO → 89.31% ± 10.38%"),
        ("2", "Geo × YOLO 协同效应 (最大发现)",
         "YOLO wider crop 保留空间上下文, Geometric 训练出姿态不变性, 增益 ×7.8"),
        ("3", "小数据集 FER 铁律",
         "正则化全部有害, ColorJitter 灾难性 → 数据增强需谨慎选择"),
    ]

    for i, (num, title, desc) in enumerate(conclusions):
        y = Inches(1.5) + i * Inches(1.6)

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_TABLE_HEADER_BG
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]
        cp.text = num
        cp.font.size = Pt(20)
        cp.font.color.rgb = C_WHITE
        cp.font.bold = True
        cp.alignment = PP_ALIGN.CENTER

        # Title
        txb1 = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.0), Inches(10.0), Inches(0.4))
        tp1 = txb1.text_frame.paragraphs[0]
        tp1.text = title
        tp1.font.size = Pt(20)
        tp1.font.bold = True
        tp1.font.color.rgb = C_TITLE

        # Description
        txb2 = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.45), Inches(10.0), Inches(0.5))
        tp2 = txb2.text_frame.paragraphs[0]
        tp2.text = desc
        tp2.font.size = Pt(14)
        tp2.font.color.rgb = C_SUBTLE

    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(6.0), Inches(11.0), Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    line.line.fill.background()

    # Next steps
    txb3 = slide.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(3.0), Inches(0.4))
    tp3 = txb3.text_frame.paragraphs[0]
    tp3.text = "下一步方向"
    tp3.font.size = Pt(14)
    tp3.font.bold = True
    tp3.font.color.rgb = C_TITLE

    next_steps = "更大 Backbone (ResNet-101, ViT)  |  跨数据集验证 (RAF-DB, AffectNet)  |  论文对标: 89.31% 远超已有方法 79%"
    txb4 = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.0), Inches(0.4))
    tp4 = txb4.text_frame.paragraphs[0]
    tp4.text = next_steps
    tp4.font.size = Pt(11)
    tp4.font.color.rgb = C_SUBTLE


# ═══════════════════════════════════════════
# Main
# ═══════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_slide_1_cover(prs)
    add_slide_2_overview(prs)
    add_slide_3_main_table(prs)
    add_slide_4_backbone(prs)
    add_slide_5_optimizer(prs)
    add_slide_6_reg_aug(prs)
    add_slide_7_synergy(prs)
    add_slide_8_f1_folds(prs)
    add_slide_9_summary_table(prs)
    add_slide_10_conclusion(prs)

    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, 'results', '阶段二', 'Phase2-组会汇报.pptx')
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")


if __name__ == '__main__':
    main()
